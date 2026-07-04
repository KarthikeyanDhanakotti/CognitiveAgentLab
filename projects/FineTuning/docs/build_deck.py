"""
Build Week9_FineTuning_Domain_Adaptation.pptx from scratch.

Runs standalone — no dependency on the previous deck file.
Content is derived from the workshop scripts under this repo.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Design tokens ───────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0F, 0x1E, 0x3E)   # deep navy — headers, bars
TEAL       = RGBColor(0x00, 0x8C, 0x9E)   # secondary accent
ORANGE     = RGBColor(0xE8, 0x71, 0x1C)   # highlight, "before" state
GREEN      = RGBColor(0x2E, 0x8B, 0x57)   # "after" / success
RED        = RGBColor(0xC0, 0x39, 0x2B)   # regression / warning
GREY_DARK  = RGBColor(0x3A, 0x3A, 0x3A)
GREY_MED   = RGBColor(0x6E, 0x6E, 0x6E)
GREY_LIGHT = RGBColor(0xEE, 0xEE, 0xEE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD  = "Segoe UI Semibold"
FONT_BODY  = "Segoe UI"
FONT_MONO  = "Consolas"

# 16:9 slide size (widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, left, top, width, height, text,
                *, font=FONT_BODY, size=18, bold=False, color=GREY_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items,
                *, size=18, color=GREY_DARK, bullet="•", spacing=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = f"{bullet}  {item}" if bullet else item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_slide_frame(prs, *, section_label=None, title=None, subtitle=None,
                    footer_num=None, footer_total=None):
    """A blank slide with our standard header bar + optional title/subtitle."""
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # Left navy sidebar
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, NAVY)

    # Section label (top-left small tag)
    if section_label:
        add_textbox(slide, Inches(0.6), Inches(0.35), Inches(6), Inches(0.35),
                    section_label.upper(),
                    font=FONT_HEAD, size=12, color=TEAL, bold=True)

    # Title
    if title:
        add_textbox(slide, Inches(0.6), Inches(0.65), Inches(12), Inches(0.9),
                    title, font=FONT_HEAD, size=32, bold=True, color=NAVY)

    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
                    subtitle, font=FONT_BODY, size=16, color=GREY_MED)

    # Footer
    add_textbox(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.35),
                "Week 9 · Fine-Tuning & Domain Adaptation · Karthikeyan Dhanakotti",
                font=FONT_BODY, size=10, color=GREY_MED)
    if footer_num and footer_total:
        add_textbox(slide, Inches(12.3), Inches(7.05), Inches(1), Inches(0.35),
                    f"{footer_num} / {footer_total}",
                    font=FONT_BODY, size=10, color=GREY_MED,
                    align=PP_ALIGN.RIGHT)
    return slide


# ── Slide builders ──────────────────────────────────────────────────────────
def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Teal accent bar
    add_rect(slide, 0, Inches(4.6), SLIDE_W, Inches(0.08), TEAL)

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5),
                "WEEK 9  ·  COGNITIVE AGENT LAB",
                font=FONT_HEAD, size=14, color=TEAL, bold=True)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.6),
                "Fine-Tuning &\nDomain Adaptation",
                font=FONT_HEAD, size=54, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.8), Inches(4.85), Inches(11.5), Inches(0.5),
                "Building a Healthcare LLM Assistant with QLoRA · HF Hub · LangSmith",
                font=FONT_BODY, size=20, color=WHITE)

    add_textbox(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4),
                "Karthikeyan Dhanakotti   |   4-hour hands-on webinar   |   Runs on Google Colab (free T4)",
                font=FONT_BODY, size=12, color=GREY_LIGHT)


def slide_welcome(prs):
    slide = add_slide_frame(prs, section_label="Welcome",
                            title="What you'll build today",
                            subtitle="One end-to-end production loop in under 4 hours")
    add_bullets(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4.2), [
        "A healthcare Q&A assistant fine-tuned with QLoRA on Qwen2.5-1.5B-Instruct.",
        "Deployed to Hugging Face Hub as a ~15 MB LoRA adapter.",
        "Evaluated with GPT-4o-mini as an LLM-as-Judge on Helpfulness / Accuracy / Safety.",
        "All of it runnable on a free Google Colab T4 GPU.",
        "Total out-of-pocket cost: ≈ $1–2.",
    ], size=20, spacing=10)

    # Callout box
    add_rect(slide, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.1), GREY_LIGHT)
    add_textbox(slide, Inches(1.1), Inches(5.85), Inches(11.1), Inches(0.9),
                "The most important skill you'll leave with: knowing when NOT to fine-tune — "
                "and how to prove it works when you do.",
                font=FONT_HEAD, size=16, color=NAVY, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


def slide_agenda(prs):
    slide = add_slide_frame(prs, section_label="Agenda",
                            title="4-hour running order")

    rows = [
        ("0:00",  "10 min", "Welcome + verify accounts",                     GREY_MED),
        ("0:10",  "50 min", "Module 1 — Strategy & Data",                    TEAL),
        ("1:00",  "10 min", "☕ Break",                                       GREY_MED),
        ("1:10",  "70 min", "Module 2 — QLoRA Fine-Tuning (live demo)",      TEAL),
        ("2:20",  "10 min", "☕ Break",                                       GREY_MED),
        ("2:30",  "30 min", "Module 3 — Deployment & Inference",             TEAL),
        ("3:00",  "50 min", "Module 4 — LangSmith Evaluation",               TEAL),
        ("3:50",  "10 min", "Wrap-up + Q&A",                                 GREY_MED),
    ]

    top = Inches(2.2)
    for i, (t, dur, label, color) in enumerate(rows):
        y = top + Inches(0.55 * i)
        add_rect(slide, Inches(0.9), y, Inches(0.15), Inches(0.45), color)
        add_textbox(slide, Inches(1.2), y, Inches(1.2), Inches(0.45),
                    t, font=FONT_HEAD, size=16, color=NAVY, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(2.4), y, Inches(1.4), Inches(0.45),
                    dur, font=FONT_BODY, size=14, color=GREY_MED,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(3.8), y, Inches(9), Inches(0.45),
                    label, font=FONT_BODY, size=16, color=GREY_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)


def slide_prereqs(prs):
    slide = add_slide_frame(prs, section_label="Pre-flight",
                            title="Accounts you need — verify NOW",
                            subtitle="The #1 blocker is the wrong Hugging Face token type")

    headers = ["Account", "Used for", "Setup", "Cost"]
    rows = [
        ("Google Colab",         "Modules 2 & 3",         "Already have",  "Free"),
        ("Hugging Face (Write)", "Modules 2 & 3",         "2 min",         "Free"),
        ("OpenAI API key",       "Modules 2 (prep), 4",   "5 min",         "~$1–2"),
        ("LangSmith",            "Module 4",              "3 min",         "Free"),
    ]

    left = Inches(0.9); top = Inches(2.4)
    col_w = [Inches(3.2), Inches(3.6), Inches(2.5), Inches(2.5)]

    # Header row
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.15), top, col_w[i], Inches(0.55),
                    h, font=FONT_HEAD, size=14, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    # Body rows
    for r, row in enumerate(rows):
        y = top + Inches(0.55 * (r + 1))
        x = left
        row_fill = WHITE if r % 2 == 0 else GREY_LIGHT
        for i, cell in enumerate(row):
            add_rect(slide, x, y, col_w[i], Inches(0.5), row_fill)
            add_textbox(slide, x + Inches(0.15), y, col_w[i], Inches(0.5),
                        cell, font=FONT_BODY, size=14, color=GREY_DARK,
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_rect(slide, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9), ORANGE)
    add_textbox(slide, Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.9),
                "Reminder — the HF token MUST be a Write token, not Read. "
                "This trips up ~15% of attendees every time.",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


# ── Module 1 ────────────────────────────────────────────────────────────────
def slide_m1_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.06), TEAL)

    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
                "MODULE 1  ·  50 MIN  ·  NO GPU",
                font=FONT_HEAD, size=14, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.5), Inches(1.2),
                "Strategy & Dataset Preparation",
                font=FONT_HEAD, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.6),
                "When to fine-tune, and what data actually works.",
                font=FONT_BODY, size=20, color=GREY_LIGHT)


def slide_three_approaches(prs):
    slide = add_slide_frame(prs, section_label="Module 1",
                            title="Three approaches to domain adaptation",
                            subtitle="Each step is more expensive and less reversible than the last")

    boxes = [
        ("Prompt engineering", "Change WHAT you send",
         ["$0  (tokens only)", "Hours to try", "Reversible", "Solves 80% of cases"],
         GREEN),
        ("RAG",                "Change what CONTEXT the model sees",
         ["Vector DB + embed", "Days to weeks", "Reversible", "Best for factual grounding"],
         TEAL),
        ("Fine-tuning",        "Change the MODEL itself",
         ["GPU compute + data", "Days to weeks", "NOT reversible", "Best for style / vocabulary"],
         ORANGE),
    ]

    top = Inches(2.3); left = Inches(0.9)
    col_w = Inches(3.9); gap = Inches(0.15)

    for i, (title, sub, bullets, color) in enumerate(boxes):
        x = left + (col_w + gap) * i
        add_rect(slide, x, top, col_w, Inches(0.7), color)
        add_textbox(slide, x + Inches(0.2), top, col_w, Inches(0.7),
                    title, font=FONT_HEAD, size=20, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, x, top + Inches(0.7), col_w, Inches(3.7), GREY_LIGHT)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.8), col_w - Inches(0.4),
                    Inches(0.5), sub, font=FONT_BODY, size=13, color=GREY_MED,
                    bold=True)
        add_bullets(slide, x + Inches(0.2), top + Inches(1.35),
                    col_w - Inches(0.4), Inches(3),
                    bullets, size=14, spacing=6)

    add_rect(slide, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6), NAVY)
    add_textbox(slide, Inches(1.0), Inches(6.35), Inches(11.4), Inches(0.6),
                "Golden rule:  Prompt engineering first  →  RAG second  →  Fine-tuning third.",
                font=FONT_HEAD, size=16, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


def slide_when_ft_wins(prs):
    slide = add_slide_frame(prs, section_label="Module 1",
                            title="Fine-tuning wins vs loses",
                            subtitle="If none of the left column matches, use prompting or RAG instead")

    top = Inches(2.3); w = Inches(5.6); h = Inches(4.3)

    # Wins panel
    add_rect(slide, Inches(0.9), top, w, Inches(0.55), GREEN)
    add_textbox(slide, Inches(1.1), top, w, Inches(0.55),
                "Fine-tuning WINS when…",
                font=FONT_HEAD, size=17, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.9), top + Inches(0.55), w, Inches(3.75), GREY_LIGHT)
    add_bullets(slide, Inches(1.1), top + Inches(0.7), w - Inches(0.3),
                Inches(3.5), [
                    "You need consistent output format / tone baked into weights.",
                    "Domain-specific vocabulary usage (e.g. clinical shorthand).",
                    "You want shorter prompts → lower token cost at scale.",
                    "Latency-sensitive apps (less context to process).",
                    "You need behavior that survives system-prompt tampering.",
                ], size=15, spacing=6)

    # Loses panel
    add_rect(slide, Inches(6.75), top, w, Inches(0.55), RED)
    add_textbox(slide, Inches(6.95), top, w, Inches(0.55),
                "Fine-tuning is the WRONG tool when…",
                font=FONT_HEAD, size=17, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(6.75), top + Inches(0.55), w, Inches(3.75), GREY_LIGHT)
    add_bullets(slide, Inches(6.95), top + Inches(0.7), w - Inches(0.3),
                Inches(3.5), [
                    "You need up-to-date info  →  use RAG.",
                    "You need source citations  →  use RAG.",
                    "Simple JSON / formatting  →  prompt engineering.",
                    "<100 queries/day  →  prompt engineering.",
                    "You're chasing new knowledge — the model already knows it.",
                ], size=15, spacing=6)


def slide_qlora_5min(prs):
    slide = add_slide_frame(prs, section_label="Module 1",
                            title="QLoRA in 5 minutes",
                            subtitle="Why this webinar is possible on a free Colab GPU")

    stages = [
        ("Full Fine-Tuning", "Update ALL 1.5B params",
         "Needs A100 · 40 GB+ VRAM", RED),
        ("LoRA",             "Freeze weights. Train small side-paths",
         "Only ~9M params · ~6 GB", ORANGE),
        ("QLoRA",            "Same as LoRA + compress base to 4-bit",
         "3 GB → 0.75 GB · fits on T4", GREEN),
    ]

    top = Inches(2.4); left = Inches(0.9)
    col_w = Inches(3.9); gap = Inches(0.15)

    for i, (name, what, cost, color) in enumerate(stages):
        x = left + (col_w + gap) * i
        add_rect(slide, x, top, col_w, Inches(3.4), GREY_LIGHT)
        add_rect(slide, x, top, col_w, Inches(0.15), color)

        add_textbox(slide, x + Inches(0.2), top + Inches(0.3), col_w, Inches(0.6),
                    name, font=FONT_HEAD, size=22, color=NAVY, bold=True)
        add_textbox(slide, x + Inches(0.2), top + Inches(1.05), col_w - Inches(0.3),
                    Inches(1.2), what, font=FONT_BODY, size=15, color=GREY_DARK)
        add_rect(slide, x + Inches(0.2), top + Inches(2.5),
                 col_w - Inches(0.4), Inches(0.75), color)
        add_textbox(slide, x + Inches(0.2), top + Inches(2.5),
                    col_w - Inches(0.4), Inches(0.75),
                    cost, font=FONT_HEAD, size=14, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.7),
                "Before 2023, fine-tuning a 1.5B model needed a $5K GPU.  "
                "Now Google gives you one for free.",
                font=FONT_HEAD, size=17, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER)


def slide_dataset_preview(prs):
    slide = add_slide_frame(prs, section_label="Module 1",
                            title="Two datasets, two futures",
                            subtitle="Same model, same LoRA config — only the DATA differs")

    headers = ["Metric",                  "ChatDoctor (v1)",   "WikiDoc reformatted (v2)"]
    rows = [
        ("Total examples",                "112,165",           "2,100"),
        ("Persona contamination",         "63.1 %",            "0.0 %"),
        ("Boilerplate sign-offs",         "28.2 %",            "0.0 %"),
        ("Safety disclaimers",            "3.2 %",             "99.4 %"),
        ("Avg answer length (chars)",     "603",               "910"),
    ]

    left = Inches(0.9); top = Inches(2.3)
    col_w = [Inches(4.3), Inches(3.6), Inches(3.6)]

    # Header
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.15), top, col_w[i], Inches(0.55),
                    h, font=FONT_HEAD, size=14, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    for r, row in enumerate(rows):
        y = top + Inches(0.55 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.5), fill)
            color = GREY_DARK
            if i == 1 and r in (1, 2): color = RED
            if i == 2 and r in (1, 3): color = GREEN
            add_textbox(slide, x + Inches(0.15), y, col_w[i], Inches(0.5),
                        cell, font=FONT_BODY, size=14, color=color,
                        bold=(i > 0 and r in (1, 2, 3)),
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_rect(slide, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.85), NAVY)
    add_textbox(slide, Inches(1.1), Inches(6.05), Inches(11.3), Inches(0.85),
                "If 63% of your training data starts with 'Hi, welcome to Chat Doctor' — "
                "what will your model learn to say?",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


def slide_m1_handson(prs):
    slide = add_slide_frame(prs, section_label="Module 1 · Hands-on",
                            title="Dataset audit exercise (25 min)",
                            subtitle="Notebook: module_1_strategy_data/notebooks/01_dataset_quality_audit.ipynb")

    parts = [
        ("Part A",  "Load ChatDoctor. Run regex audit for persona strings and boilerplate."),
        ("Part B",  "Load reformatted WikiDoc. Run the SAME audit. Compare metrics."),
        ("Part C",  "Prompt-engineering baseline — is the base model already good enough?"),
    ]
    y = Inches(2.4)
    for tag, desc in parts:
        add_rect(slide, Inches(0.9), y, Inches(1.4), Inches(0.85), TEAL)
        add_textbox(slide, Inches(0.9), y, Inches(1.4), Inches(0.85),
                    tag, font=FONT_HEAD, size=18, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, Inches(2.3), y, Inches(10.1), Inches(0.85), GREY_LIGHT)
        add_textbox(slide, Inches(2.5), y, Inches(9.9), Inches(0.85),
                    desc, font=FONT_BODY, size=15, color=GREY_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.1)


# ── Module 2 ────────────────────────────────────────────────────────────────
def slide_m2_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.06), TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
                "MODULE 2  ·  70 MIN  ·  GPU: T4",
                font=FONT_HEAD, size=14, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.5), Inches(1.2),
                "QLoRA Fine-Tuning on Google Colab",
                font=FONT_HEAD, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.6),
                "Same model. Same config. Two datasets. Opposite outcomes.",
                font=FONT_BODY, size=20, color=GREY_LIGHT)


def slide_two_experiments(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="The two experiments",
                            subtitle="Only ONE thing changes between them — the training data")

    headers = ["",                     "v1 notebook",              "v2 notebook"]
    rows = [
        ("Base model",                 "Qwen2.5-1.5B-Instruct",    "Qwen2.5-1.5B-Instruct"),
        ("LoRA config",                "r=16, α=32, dropout=0.05", "r=16, α=32, dropout=0.05"),
        ("Optimizer / batch",          "Identical",                "Identical"),
        ("Dataset",                    "ChatDoctor (noisy, 112k)", "WikiDoc (clean, 2.1k)"),
        ("Predicted outcome",          "Model gets WORSE",         "Model gets BETTER"),
    ]

    left = Inches(0.9); top = Inches(2.3)
    col_w = [Inches(3.5), Inches(4.0), Inches(4.0)]
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.15), top, col_w[i], Inches(0.55),
                    h, font=FONT_HEAD, size=14, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]

    for r, row in enumerate(rows):
        y = top + Inches(0.55 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.5), fill)
            color = GREY_DARK; bold = False
            if r == 4 and i == 1: color = RED; bold = True
            if r == 4 and i == 2: color = GREEN; bold = True
            if r == 3 and i > 0: bold = True
            add_textbox(slide, x + Inches(0.15), y, col_w[i], Inches(0.5),
                        cell, font=FONT_BODY, size=14, color=color, bold=bold,
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_textbox(slide, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.7),
                "Data quality — not model choice, not hyperparameter magic — "
                "is what makes or breaks a fine-tune.",
                font=FONT_HEAD, size=15, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER)


def slide_pipeline(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="Training pipeline at a glance")

    steps = [
        ("1", "GPU check", "!nvidia-smi"),
        ("2", "Install libs", "transformers · peft · trl · bitsandbytes"),
        ("3", "Load base model", "Qwen2.5-1.5B in 4-bit"),
        ("4", "BEFORE benchmark", "10 medical prompts → baseline answers"),
        ("5", "Attach LoRA", "Only ~0.6% of params trainable"),
        ("6", "Train", "500 steps · ~20–30 min on T4"),
        ("7", "AFTER benchmark", "Same 10 prompts → compare"),
        ("8", "Push to HF Hub", "~15 MB adapter, usable from anywhere"),
    ]

    top = Inches(2.2); left = Inches(0.9)
    for i, (n, title, desc) in enumerate(steps):
        col = i % 4; row = i // 4
        x = left + Inches(col * 3.05)
        y = top + Inches(row * 2.15)

        add_rect(slide, x, y, Inches(2.85), Inches(1.9), GREY_LIGHT)
        add_rect(slide, x, y, Inches(2.85), Inches(0.35), NAVY)
        add_textbox(slide, x + Inches(0.2), y, Inches(0.6), Inches(0.35),
                    f"STEP {n}", font=FONT_HEAD, size=11, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(2.5),
                    Inches(0.55), title, font=FONT_HEAD, size=15, color=NAVY,
                    bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.05), Inches(2.5),
                    Inches(0.75), desc, font=FONT_BODY, size=11, color=GREY_MED)


def slide_quantization(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="4-bit quantization with BitsAndBytes",
                            subtitle="Why a 1.5B model fits comfortably on a 15 GB T4")

    code = ('BitsAndBytesConfig(\n'
            '    load_in_4bit=True,\n'
            '    bnb_4bit_quant_type="nf4",              # NormalFloat 4-bit\n'
            '    bnb_4bit_compute_dtype=torch.bfloat16,  # compute in bf16\n'
            '    bnb_4bit_use_double_quant=True,         # quantize the constants too\n'
            ')')

    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(7), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = FONT_MONO; run.font.size = Pt(15)
        run.font.color.rgb = NAVY

    # Model-size infographic
    right_x = Inches(8.4); top = Inches(2.3); w = Inches(4.2)
    add_rect(slide, right_x, top, w, Inches(3.4), GREY_LIGHT)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.15), w, Inches(0.4),
                "Model size on GPU",
                font=FONT_HEAD, size=14, color=NAVY, bold=True)

    # fp16 bar
    add_rect(slide, right_x + Inches(0.2), top + Inches(0.8),
             Inches(3.8), Inches(0.5), ORANGE)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.8),
                Inches(3.8), Inches(0.5),
                "fp16  ·  3.0 GB",
                font=FONT_HEAD, size=13, color=WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 4-bit bar (proportional)
    add_rect(slide, right_x + Inches(0.2), top + Inches(1.7),
             Inches(0.95), Inches(0.5), GREEN)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(1.7),
                Inches(3.8), Inches(0.5),
                "4-bit  ·  0.75 GB",
                font=FONT_HEAD, size=13, color=NAVY, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, right_x + Inches(0.2), top + Inches(2.55),
                w - Inches(0.4), Inches(0.7),
                "73% VRAM reduction with negligible quality loss.\n"
                "Weights stored 4-bit; computation stays in bf16.",
                font=FONT_BODY, size=13, color=GREY_DARK)

    add_rect(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.4), GREY_LIGHT)
    add_textbox(slide, Inches(1.1), Inches(5.55), Inches(11.1), Inches(1.2),
                "Rule at inference: use the EXACT same quantization config as during training. "
                "The LoRA adapter was calibrated against 4-bit weights — load base in fp16 "
                "and adapter corrections are aimed at the wrong numbers.",
                font=FONT_BODY, size=15, color=GREY_DARK)


def slide_lora_config(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="LoRA config — what each knob does")

    code = ('LoraConfig(\n'
            '    r=16,                        # rank of the side-path\n'
            '    lora_alpha=32,               # effective LR ≈ alpha/r = 2.0\n'
            '    lora_dropout=0.05,           # 5% dropout on adapter neurons\n'
            '    target_modules="all-linear", # attach to every linear layer\n'
            ')')
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(7), Inches(2.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = FONT_MONO; run.font.size = Pt(15)
        run.font.color.rgb = NAVY

    # Right-side analogy
    right_x = Inches(8.4); top = Inches(2.3)
    add_rect(slide, right_x, top, Inches(4.2), Inches(2.4), GREY_LIGHT)
    add_textbox(slide, right_x + Inches(0.2), top + Inches(0.15),
                Inches(3.8), Inches(2.1),
                ["Analogy",
                 "",
                 "Frozen base = a 4-lane highway.",
                 "LoRA = a small side-road that",
                 "nudges the traffic.",
                 "",
                 "r  = width of side-road",
                 "α  = amplification of its output",
                 "dropout = randomly closes 5% of lanes"],
                font=FONT_BODY, size=13, color=GREY_DARK)

    # Training math box
    add_rect(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.85), NAVY)
    add_textbox(slide, Inches(1.1), Inches(5.05), Inches(11.1), Inches(0.5),
                "Trainable-parameter math",
                font=FONT_HEAD, size=15, color=TEAL, bold=True)
    add_textbox(slide, Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.3),
                ["Trainable params:  9,232,384",
                 "All params:            1,552,946,688",
                 "Trainable %:          0.5945%   ← this is why LoRA is fast"],
                font=FONT_MONO, size=15, color=WHITE)


def slide_hyperparams(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="Training hyperparameters — do the math live")

    left_code = ("Dataset:          2,000 examples\n"
                 "Effective batch:  8   (batch=1 × grad_accum=8)\n"
                 "Steps / epoch:    2000 / 8 = 250\n"
                 "\n"
                 "v2:               250 × 2 epochs = 500 steps\n"
                 "Warmup:           50 steps = 10%")
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(6.5), Inches(3.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(left_code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = FONT_MONO; run.font.size = Pt(16)
        run.font.color.rgb = NAVY

    # Right: hyperparam comparison
    right_x = Inches(7.6); top = Inches(2.3); w = Inches(4.9)
    headers = ["Param",           "v1",       "v2"]
    rows = [
        ("Learning rate",         "2e-4",     "5e-5"),
        ("Epochs",                "3",        "2"),
        ("Warmup steps",          "75 (10%)", "50 (10%)"),
        ("Dropout",               "0.05",     "0.05"),
    ]
    col_w = [Inches(2.0), Inches(1.45), Inches(1.45)]
    x = right_x
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.5), NAVY)
        add_textbox(slide, x + Inches(0.1), top, col_w[i], Inches(0.5),
                    h, font=FONT_HEAD, size=13, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = top + Inches(0.5 * (r + 1))
        x = right_x
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.45), fill)
            add_textbox(slide, x + Inches(0.1), y, col_w[i], Inches(0.45),
                        cell, font=FONT_BODY, size=13, color=GREY_DARK,
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_textbox(slide, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.9),
                "v2 uses a gentler recipe (lower LR, fewer epochs) — "
                "the goal is to shift STYLE, not overwrite KNOWLEDGE.",
                font=FONT_HEAD, size=16, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER)


def slide_v1_failure(prs):
    slide = add_slide_frame(prs, section_label="Module 2  ·  v1",
                            title="v1 — how NOT to fine-tune",
                            subtitle="Same model. Dirty ChatDoctor data. Output regresses.")

    # BEFORE box
    top = Inches(2.3); w = Inches(5.6); h = Inches(2.2)
    add_rect(slide, Inches(0.9), top, w, Inches(0.5), GREY_DARK)
    add_textbox(slide, Inches(1.1), top, w, Inches(0.5),
                "BEFORE (base model)",
                font=FONT_HEAD, size=14, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.9), top + Inches(0.5), w, h, GREY_LIGHT)
    add_textbox(slide, Inches(1.1), top + Inches(0.6), w - Inches(0.3), h,
                "Type 2 diabetes symptoms include increased thirst, "
                "frequent urination, blurred vision, fatigue, and slow-healing "
                "wounds. These symptoms develop gradually. Please consult a "
                "healthcare professional for proper diagnosis.",
                font=FONT_BODY, size=13, color=GREY_DARK)

    # AFTER box
    add_rect(slide, Inches(6.75), top, w, Inches(0.5), RED)
    add_textbox(slide, Inches(6.95), top, w, Inches(0.5),
                "AFTER v1 (ChatDoctor fine-tuned)",
                font=FONT_HEAD, size=14, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(6.75), top + Inches(0.5), w, h, GREY_LIGHT)
    add_textbox(slide, Inches(6.95), top + Inches(0.6), w - Inches(0.3), h,
                "Hello, thanks for posting your query. The symptoms of diabetes "
                "include increased urination and thirst. Hope this helps. "
                "Wishing you good health.",
                font=FONT_BODY, size=13, color=GREY_DARK)

    add_rect(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.55), NAVY)
    add_textbox(slide, Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.4),
                ["Three failure modes v1 exhibits:",
                 "1.  Persona contamination — greets like a chatbot.",
                 "2.  Lost structure — bullet points replaced by paragraph.",
                 "3.  Weakened safety — disclaimer replaced with a sign-off."],
                font=FONT_BODY, size=15, color=WHITE)


def slide_v2_success(prs):
    slide = add_slide_frame(prs, section_label="Module 2  ·  v2",
                            title="v2 — how TO fine-tune",
                            subtitle="Same model. Clean WikiDoc data. Output improves in the ways we want.")

    top = Inches(2.3); w = Inches(5.6); h = Inches(2.6)

    add_rect(slide, Inches(0.9), top, w, Inches(0.5), GREY_DARK)
    add_textbox(slide, Inches(1.1), top, w, Inches(0.5),
                "BEFORE (base model)",
                font=FONT_HEAD, size=14, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.9), top + Inches(0.5), w, h, GREY_LIGHT)
    add_textbox(slide, Inches(1.1), top + Inches(0.6), w - Inches(0.3), h,
                "Type 2 diabetes symptoms include increased thirst, "
                "frequent urination, blurred vision, fatigue and slow-healing "
                "wounds. These symptoms develop gradually.",
                font=FONT_BODY, size=13, color=GREY_DARK)

    add_rect(slide, Inches(6.75), top, w, Inches(0.5), GREEN)
    add_textbox(slide, Inches(6.95), top, w, Inches(0.5),
                "AFTER v2 (WikiDoc fine-tuned)",
                font=FONT_HEAD, size=14, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(6.75), top + Inches(0.5), w, h, GREY_LIGHT)
    add_textbox(slide, Inches(6.95), top + Inches(0.6), w - Inches(0.3), h,
                ["Common symptoms of Type 2 diabetes:",
                 "• Increased thirst (polydipsia) and frequent urination",
                 "• Unexplained fatigue and blurred vision",
                 "• Slow-healing wounds and frequent infections",
                 "• Numbness or tingling in hands and feet",
                 "",
                 "Please consult a healthcare professional for personalized advice."],
                font=FONT_BODY, size=12, color=GREY_DARK)

    add_textbox(slide, Inches(0.9), Inches(6.05), Inches(11.5), Inches(0.9),
                "Same model. Same LoRA config. Structure preserved, safety added, "
                "zero persona contamination.",
                font=FONT_HEAD, size=16, color=NAVY, bold=True,
                align=PP_ALIGN.CENTER)


def slide_push_to_hub(prs):
    slide = add_slide_frame(prs, section_label="Module 2",
                            title="Push the adapter to Hugging Face Hub",
                            subtitle="What actually lives in your HF repo")

    code = ('# One line publishes your adapter to the internet:\n'
            'model.push_to_hub("your-username/healthcare-assistant-lora-v2")\n'
            '\n'
            '# The base model (Qwen 2.5 1.5B, ~3 GB) is NOT re-uploaded —\n'
            '# it stays where Qwen already hosts it.')
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = FONT_MONO; run.font.size = Pt(15)
        run.font.color.rgb = NAVY

    add_rect(slide, Inches(0.9), Inches(4.4), Inches(11.5), Inches(2.4), GREY_LIGHT)
    add_textbox(slide, Inches(1.1), Inches(4.55), Inches(11.1), Inches(0.4),
                "your-username/healthcare-assistant-lora-v2   (~20–50 MB total)",
                font=FONT_HEAD, size=15, color=NAVY, bold=True)
    add_textbox(slide, Inches(1.1), Inches(5.0), Inches(11.1), Inches(1.9),
                ["adapter_config.json           # LoRA config (r, alpha, target modules)",
                 "adapter_model.safetensors     # trained LoRA weights (~15 MB)",
                 "tokenizer.json                # tokenizer vocabulary",
                 "tokenizer_config.json         # tokenizer settings",
                 "special_tokens_map.json       # special token mappings"],
                font=FONT_MONO, size=13, color=GREY_DARK)


# ── Module 3 ────────────────────────────────────────────────────────────────
def slide_m3_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.06), TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
                "MODULE 3  ·  30 MIN  ·  GPU: T4",
                font=FONT_HEAD, size=14, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.5), Inches(1.2),
                "Deployment & Inference",
                font=FONT_HEAD, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.6),
                "Adapters are patches. Load them, toggle them, ship them.",
                font=FONT_BODY, size=20, color=GREY_LIGHT)


def slide_adapter_toggle(prs):
    slide = add_slide_frame(prs, section_label="Module 3",
                            title="The adapter-toggle pattern",
                            subtitle="Two models in one load — save ~1 GB VRAM vs loading both")

    code = ('# Fine-tuned behavior (adapter ON):\n'
            'ft_model.enable_adapter_layers()\n'
            'finetuned_response = generate(ft_model, tokenizer, prompt)\n'
            '\n'
            '# Base model behavior (adapter OFF) — same model object!\n'
            'ft_model.disable_adapter_layers()\n'
            'base_response = generate(ft_model, tokenizer, prompt)')
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.name = FONT_MONO; run.font.size = Pt(15)
        run.font.color.rgb = NAVY

    add_rect(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.7), NAVY)
    add_textbox(slide, Inches(1.1), Inches(5.3), Inches(11.1), Inches(1.5),
                ["ON :   output = base_layer(x) + lora_B(lora_A(x)) × (α/r)",
                 "OFF:   output = base_layer(x)                                 ← pure base",
                 "",
                 "One model object.  Two behaviors.  No extra VRAM for the second."],
                font=FONT_MONO, size=14, color=WHITE)


def slide_gen_params(prs):
    slide = add_slide_frame(prs, section_label="Module 3",
                            title="Generation parameters for healthcare",
                            subtitle="Low temperature isn't optional in regulated domains")

    headers = ["Parameter",        "Healthcare setting",  "Why"]
    rows = [
        ("temperature",            "0.1 – 0.3",           "Consistent answers · avoids hallucination"),
        ("top_p",                  "0.5 – 0.7",           "Restrict to high-probability tokens"),
        ("max_new_tokens",         "256 – 512",           "Long enough to be thorough, short enough to focus"),
        ("do_sample",              "True (with low T)",   "Not deterministic, not chaotic either"),
        ("repetition_penalty",     "1.05 – 1.15",         "Prevents parroting the input back"),
    ]

    left = Inches(0.9); top = Inches(2.3)
    col_w = [Inches(3.0), Inches(3.2), Inches(5.3)]
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.15), top, col_w[i], Inches(0.55),
                    h, font=FONT_HEAD, size=14, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = top + Inches(0.55 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.5), fill)
            font = FONT_MONO if i == 0 else FONT_BODY
            add_textbox(slide, x + Inches(0.15), y, col_w[i], Inches(0.5),
                        cell, font=font, size=13, color=GREY_DARK,
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]


def slide_prod_strategies(prs):
    slide = add_slide_frame(prs, section_label="Module 3",
                            title="Production deployment — right-size the GPU")

    headers = ["Strategy",         "Example",                        "Cold start", "Cost model",     "Best for"]
    rows = [
        ("Serverless",             "HF Endpoints · SageMaker Serverless",
         "30–120 s",  "Pay-per-request", "< 100 req/day"),
        ("Dedicated GPU",          "SageMaker Real-time · Azure ML",
         "None",      "Pay-per-hour",    "Prod w/ SLA"),
        ("Container (vLLM/TGI)",   "ECS · EKS · EC2 with vLLM",
         "None",      "Pay-per-hour",    "Multi-model, K8s"),
        ("Model-as-a-Service",     "Bedrock · Azure OpenAI",
         "None",      "Pay-per-token",   "Zero-ops"),
    ]

    left = Inches(0.9); top = Inches(2.3)
    col_w = [Inches(2.4), Inches(3.4), Inches(1.4), Inches(1.8), Inches(2.5)]
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.55), NAVY)
        add_textbox(slide, x + Inches(0.1), top, col_w[i], Inches(0.55),
                    h, font=FONT_HEAD, size=13, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = top + Inches(0.55 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.55), fill)
            add_textbox(slide, x + Inches(0.1), y, col_w[i], Inches(0.55),
                        cell, font=FONT_BODY, size=12, color=GREY_DARK,
                        anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_rect(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.3), GREY_LIGHT)
    add_textbox(slide, Inches(1.1), Inches(5.7), Inches(11.1), Inches(1.15),
                ["• 1.5B in 4-bit runs happily on a T4 (~$0.74/hr).  An A100 (~$37/hr) wastes 98% of the GPU.",
                 "• vLLM with continuous batching = 2–4× throughput on the same GPU.",
                 "• For production, merge the adapter:  merged = ft_model.merge_and_unload()"],
                font=FONT_BODY, size=13, color=GREY_DARK)


# ── Module 4 ────────────────────────────────────────────────────────────────
def slide_m4_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.06), TEAL)
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.5),
                "MODULE 4  ·  50 MIN  ·  NO GPU",
                font=FONT_HEAD, size=14, color=TEAL, bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.85), Inches(11.5), Inches(1.2),
                "LangSmith Evaluation & Observability",
                font=FONT_HEAD, size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.6),
                "Turn 'it feels better' into a number.",
                font=FONT_BODY, size=20, color=GREY_LIGHT)


def slide_why_eval(prs):
    slide = add_slide_frame(prs, section_label="Module 4",
                            title="Why evaluation is non-negotiable")

    top = Inches(2.3); w = Inches(5.6); h = Inches(3.9)

    # Without eval
    add_rect(slide, Inches(0.9), top, w, Inches(0.55), RED)
    add_textbox(slide, Inches(1.1), top, w, Inches(0.55),
                "Without evaluation",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0.9), top + Inches(0.55), w, h, GREY_LIGHT)
    add_bullets(slide, Inches(1.1), top + Inches(0.7), w - Inches(0.3),
                Inches(3.7), [
                    "\"Does fine-tuning help?\" — unclear.",
                    "\"What improved?\" — subjective.",
                    "\"Did anything get worse?\" — unknown.",
                    "Two devs disagree — no way to resolve.",
                    "Ship it? Roll back? Argue forever.",
                ], size=15, spacing=10)

    # With eval
    add_rect(slide, Inches(6.75), top, w, Inches(0.55), GREEN)
    add_textbox(slide, Inches(6.95), top, w, Inches(0.55),
                "With LLM-as-Judge evaluation",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(6.75), top + Inches(0.55), w, h, GREY_LIGHT)
    add_bullets(slide, Inches(6.95), top + Inches(0.7), w - Inches(0.3),
                Inches(3.7), [
                    "3 objective metrics per response.",
                    "60 judge calls · ~$0.06 · 2–3 min.",
                    "Reproducible, shareable numbers.",
                    "Catches regressions before shipping.",
                    "Answers all three questions with data.",
                ], size=15, spacing=10)


def slide_three_evaluators(prs):
    slide = add_slide_frame(prs, section_label="Module 4",
                            title="The three evaluators",
                            subtitle="Each on a 0–5 scale, normalized to 0.0–1.0")

    evals = [
        ("Helpfulness",  "Is the answer actionable?",
         ["Clear bullet list + next steps",
          "Explains mechanism, not just names",
          "Directs the user to appropriate care"],
         TEAL),
        ("Accuracy",     "Is the medical content correct?",
         ["Consistent with standard-of-care",
          "No contradictions of guidelines",
          "No fabricated drug interactions"],
         NAVY),
        ("Safety",       "Does it recommend consulting a professional?",
         ["Explicit disclaimer present",
          "No definitive diagnosis given",
          "No dangerous self-treatment advice"],
         ORANGE),
    ]

    top = Inches(2.3); left = Inches(0.9)
    col_w = Inches(3.9); gap = Inches(0.15)
    for i, (name, q, bullets, color) in enumerate(evals):
        x = left + (col_w + gap) * i
        add_rect(slide, x, top, col_w, Inches(0.7), color)
        add_textbox(slide, x + Inches(0.2), top, col_w, Inches(0.7),
                    name, font=FONT_HEAD, size=22, color=WHITE, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, x, top + Inches(0.7), col_w, Inches(4), GREY_LIGHT)
        add_textbox(slide, x + Inches(0.2), top + Inches(0.85), col_w - Inches(0.4),
                    Inches(0.6), q, font=FONT_BODY, size=14, color=GREY_MED,
                    bold=True)
        add_bullets(slide, x + Inches(0.2), top + Inches(1.5),
                    col_w - Inches(0.4), Inches(3),
                    bullets, size=13, spacing=6)


def slide_eval_flow(prs):
    slide = add_slide_frame(prs, section_label="Module 4",
                            title="Evaluation pipeline",
                            subtitle="Module 2/3 produces the outputs · Module 4 scores them")

    steps = [
        ("Module 2 / 3",   "Produces inference_results.json\n10 prompts × base + fine-tuned"),
        ("Create dataset", "Upload 10 prompts as a LangSmith dataset"),
        ("Base experiment","Score base outputs — 30 judge calls (10×3 evaluators)"),
        ("FT experiment",  "Score fine-tuned outputs — 30 more calls"),
        ("Compare view",   "Side-by-side deltas in the LangSmith UI"),
    ]

    top = Inches(2.5); left = Inches(0.9); w = Inches(11.5)
    for i, (title, desc) in enumerate(steps):
        y = top + Inches(0.85 * i)
        add_rect(slide, left, y, Inches(0.75), Inches(0.75), NAVY)
        add_textbox(slide, left, y, Inches(0.75), Inches(0.75),
                    str(i + 1), font=FONT_HEAD, size=22, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, left + Inches(0.85), y, Inches(10.65), Inches(0.75),
                 GREY_LIGHT)
        add_textbox(slide, left + Inches(1.0), y, Inches(3.5), Inches(0.75),
                    title, font=FONT_HEAD, size=16, color=NAVY, bold=True,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left + Inches(4.5), y, Inches(6.9), Inches(0.75),
                    desc, font=FONT_BODY, size=13, color=GREY_DARK,
                    anchor=MSO_ANCHOR.MIDDLE)


def slide_results_v1(prs):
    slide = add_slide_frame(prs, section_label="Module 4  ·  v1 results",
                            title="v1 (ChatDoctor) — all three metrics regressed",
                            subtitle="This is what happens when you fine-tune on dirty data")

    headers = ["Metric",        "Base",  "Fine-Tuned v1",  "Delta"]
    rows = [
        ("Helpfulness",         "0.78",  "0.60",           "−0.18"),
        ("Accuracy",            "0.72",  "0.64",           "−0.08"),
        ("Safety",              "0.70",  "0.64",           "−0.06"),
    ]

    left = Inches(1.9); top = Inches(2.4)
    col_w = [Inches(2.5), Inches(2.0), Inches(2.5), Inches(2.5)]
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.6), NAVY)
        add_textbox(slide, x + Inches(0.1), top, col_w[i], Inches(0.6),
                    h, font=FONT_HEAD, size=15, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = top + Inches(0.6 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.55), fill)
            color = GREY_DARK; bold = False
            if i == 3: color = RED; bold = True
            add_textbox(slide, x + Inches(0.1), y, col_w[i], Inches(0.55),
                        cell, font=FONT_BODY, size=15, color=color, bold=bold,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_rect(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.5), RED)
    add_textbox(slide, Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.3),
                "All three metrics went DOWN. Without evaluation, someone posts "
                "'we fine-tuned a healthcare model!' on LinkedIn and ships it. "
                "With evaluation, we see: we broke the model — and we can prove it.",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


def slide_results_v2(prs):
    slide = add_slide_frame(prs, section_label="Module 4  ·  v2 results",
                            title="v2 (WikiDoc) — mixed: 2 wins + 1 regression",
                            subtitle="Clean data helped accuracy and safety, but answers got too short")

    headers = ["Metric",        "Base",  "Fine-Tuned v2",  "Delta"]
    rows = [
        ("Accuracy",            "0.66",  "0.72",           "+0.06"),
        ("Helpfulness",         "0.72",  "0.56",           "−0.16"),
        ("Safety",              "0.76",  "0.86",           "+0.10"),
    ]
    deltas_pos = {0, 2}  # rows with green delta

    left = Inches(1.9); top = Inches(2.4)
    col_w = [Inches(2.5), Inches(2.0), Inches(2.5), Inches(2.5)]
    x = left
    for i, h in enumerate(headers):
        add_rect(slide, x, top, col_w[i], Inches(0.6), NAVY)
        add_textbox(slide, x + Inches(0.1), top, col_w[i], Inches(0.6),
                    h, font=FONT_HEAD, size=15, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_w[i]
    for r, row in enumerate(rows):
        y = top + Inches(0.6 * (r + 1))
        x = left
        for i, cell in enumerate(row):
            fill = WHITE if r % 2 == 0 else GREY_LIGHT
            add_rect(slide, x, y, col_w[i], Inches(0.55), fill)
            color = GREY_DARK; bold = False
            if i == 3:
                color = GREEN if r in deltas_pos else RED
                bold = True
            add_textbox(slide, x + Inches(0.1), y, col_w[i], Inches(0.55),
                        cell, font=FONT_BODY, size=15, color=color, bold=bold,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += col_w[i]

    add_rect(slide, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.5), NAVY)
    add_textbox(slide, Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.3),
                "Two of three metrics improved. One regressed by 16 points. "
                "Without the eval we'd have shipped this thinking it was a pure win. "
                "The regression check is the whole point.",
                font=FONT_HEAD, size=15, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


# ── Closing ─────────────────────────────────────────────────────────────────
def slide_takeaways(prs):
    slide = add_slide_frame(prs, section_label="Wrap-up",
                            title="Key takeaways")

    items = [
        ("1", "Prompt engineering first · RAG second · Fine-tuning third.",
         "Each step is more expensive and less reversible than the last."),
        ("2", "Data quality determines 90% of the outcome.",
         "Same model + dirty data = worse. Same model + clean data = better."),
        ("3", "QLoRA on a free T4 is a real workflow, not a toy.",
         "1.5B model in 4-bit + 9M trainable params + 500 steps + 30 minutes."),
        ("4", "Adapters are 15 MB patches — deploy with the toggle pattern.",
         "Two behaviors in one load. Use identical quantization at inference."),
        ("5", "Fine-tuning without evaluation is opinion, not engineering.",
         "Answer the regression question — 'did anything get worse?' — every time."),
    ]

    top = Inches(2.2); left = Inches(0.9)
    for i, (n, head, sub) in enumerate(items):
        y = top + Inches(0.95 * i)
        add_rect(slide, left, y, Inches(0.75), Inches(0.75), TEAL)
        add_textbox(slide, left, y, Inches(0.75), Inches(0.75),
                    n, font=FONT_HEAD, size=22, color=WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left + Inches(1.0), y - Inches(0.05),
                    Inches(11.3), Inches(0.45),
                    head, font=FONT_HEAD, size=16, color=NAVY, bold=True)
        add_textbox(slide, left + Inches(1.0), y + Inches(0.4),
                    Inches(11.3), Inches(0.4),
                    sub, font=FONT_BODY, size=13, color=GREY_MED)


def slide_next_steps(prs):
    slide = add_slide_frame(prs, section_label="Wrap-up",
                            title="Next steps",
                            subtitle="Take what you built today into your own domain")

    add_bullets(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(4), [
        "Swap the dataset — legal, financial, customer support — the pipeline is identical.",
        "Try a bigger base model (Qwen 3B · Llama 3.2 3B · Mistral 7B in 4-bit).",
        "Add production tracing with LangSmith once you have real user traffic.",
        "Read notes.md in each module for deep-dive references we skipped live.",
        "Fork the repo — bring your own domain and share your results.",
    ], size=18, spacing=10)

    add_rect(slide, Inches(0.9), Inches(5.8), Inches(11.5), Inches(1.1), NAVY)
    add_textbox(slide, Inches(1.1), Inches(5.9), Inches(11.1), Inches(0.9),
                "github.com/KarthikeyanDhanakotti/CognitiveAgentLab/tree/main/projects/FineTuning",
                font=FONT_MONO, size=17, color=WHITE, bold=True,
                anchor=MSO_ANCHOR.MIDDLE)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, 0, Inches(4.6), SLIDE_W, Inches(0.08), TEAL)

    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.4),
                "Thank you.",
                font=FONT_HEAD, size=64, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.7),
                "Questions?",
                font=FONT_HEAD, size=28, color=TEAL, bold=True)

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.4),
                "Karthikeyan Dhanakotti  ·  Senior Data Science Leader, Microsoft",
                font=FONT_BODY, size=16, color=GREY_LIGHT)
    add_textbox(slide, Inches(0.8), Inches(5.45), Inches(11.5), Inches(0.4),
                "github.com/KarthikeyanDhanakotti/CognitiveAgentLab",
                font=FONT_MONO, size=14, color=GREY_LIGHT)


# ── Assemble ────────────────────────────────────────────────────────────────
def build(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,
        slide_welcome,
        slide_agenda,
        slide_prereqs,

        slide_m1_divider,
        slide_three_approaches,
        slide_when_ft_wins,
        slide_qlora_5min,
        slide_dataset_preview,
        slide_m1_handson,

        slide_m2_divider,
        slide_two_experiments,
        slide_pipeline,
        slide_quantization,
        slide_lora_config,
        slide_hyperparams,
        slide_v1_failure,
        slide_v2_success,
        slide_push_to_hub,

        slide_m3_divider,
        slide_adapter_toggle,
        slide_gen_params,
        slide_prod_strategies,

        slide_m4_divider,
        slide_why_eval,
        slide_three_evaluators,
        slide_eval_flow,
        slide_results_v1,
        slide_results_v2,

        slide_takeaways,
        slide_next_steps,
        slide_thanks,
    ]
    for b in builders:
        b(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return len(builders)


if __name__ == "__main__":
    out = Path(r"C:\Users\kartdh\repos\CognitiveAgentLab\projects\FineTuning"
               r"\docs\Week9_FineTuning_Domain_Adaptation.pptx")
    n = build(out)
    print(f"OK — {n} slides written to {out}")
    print(f"Size: {out.stat().st_size / 1024:.1f} KB")
